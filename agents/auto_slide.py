import google.generativeai as genai
import os
import sys
import re
import time
from pptx import Presentation
from pathlib import Path
from dotenv import load_dotenv

# 載入環境變數
load_dotenv(dotenv_path=r"C:\Users\user\claude AI_Agent\.env")

# 1. 設定你的 API Key
api_key = os.getenv("GOOGLE_API_KEY") or os.getenv("GEMINI_API_KEY")
if not api_key:
    print("錯誤: 未設定 GOOGLE_API_KEY 或 GEMINI_API_KEY 環境變數。")
    sys.exit(1)

genai.configure(api_key=api_key)

def send_notification(file_path, output_pptx, content):
    """
    透過電子郵件通知投影片已完成，並包含內容摘要
    """
    try:
        import importlib.util as _ilu
        BASE_DIR = Path(r"C:\Users\user\claude AI_Agent")
        spec = _ilu.spec_from_file_location("email_notify", str(BASE_DIR / "agents" / "email_notify.py"))
        email_mod = _ilu.module_from_spec(spec)
        spec.loader.exec_module(email_mod)
        
        # 將內容中的換行轉換為 HTML 換行
        formatted_content = content.replace("\n", "<br>")
        
        subject = f"[Yisheng AI] 媒體總結投影片已完成 - {Path(file_path).name}"
        body = f"""
        <h3>🎬 媒體總結報告已產出</h3>
        <p><b>原始檔案：</b> {file_path}</p>
        <p><b>產出投影片：</b> {output_pptx}</p>
        <hr>
        <h4>內容摘要：</h4>
        <div style="background:#f4f4f4; padding:15px; border-radius:5px; font-family:serif;">
            {formatted_content}
        </div>
        <hr>
        <p style="color:#888;font-size:12px;">此郵件由 Yisheng AI 自動發送</p>
        """
        email_mod.send_email(subject, body)
        print(f"電子郵件通知已發送。")
    except Exception as e:
        print(f"發送電子郵件通知失敗: {e}")

def process_media_to_slides(file_path):
    file_path = Path(file_path)
    if not file_path.exists():
        print(f"錯誤: 找不到檔案 {file_path}")
        return

    # 使用您清單中明確存在的模型別名
    model_name = 'gemini-flash-latest'
    print(f"正在啟動模型: {model_name}...")
    
    try:
        model = genai.GenerativeModel(model_name)
    except Exception as e:
        print(f"模型啟動失敗: {e}")
        return
    
    # 2. 上傳影片/錄音
    print(f"正在上傳並分析檔案: {file_path}...")
    try:
        media_file = genai.upload_file(path=str(file_path))
        print(f"檔案已上傳，正在等待 Google 處理 (ID: {media_file.name})...")
        
        # 等待檔案進入 ACTIVE 狀態
        while media_file.state.name == "PROCESSING":
            print(".", end="", flush=True)
            time.sleep(5)
            media_file = genai.get_file(media_file.name)
            
        if media_file.state.name != "ACTIVE":
            print(f"\n檔案處理異常: {media_file.state.name}")
            return
        
        print(f"\n檔案已準備就緒，開始生成投影片內容...")
        
    except Exception as e:
        print(f"\n上傳或處理失敗: {e}")
        return
    
    # 3. 讓 Gemini 生成投影片內容
    prompt = """
    請分析這段影片/錄音，並將重點整理成 5 頁投影片。
    每頁請提供：標題(Title) 和 3 個內文要點(Bullet Points)。
    請直接以繁體中文回答。
    輸出的格式請嚴格遵守：
    第 X 頁：[標題]
    - 要點 1
    - 要點 2
    - 要點 3
    """
    
    try:
        response = model.generate_content([prompt, media_file])
        content = response.text
    except Exception as e:
        print(f"生成內容失敗: {e}")
        return
    
    # 4. 建立實體投影片檔案
    prs = Presentation()
    pages = re.split(r"第\s*[0-9一二三四五六七八九十]+\s*頁[:：]", content)
    
    slide_count = 0
    for page in pages:
        page = page.strip()
        if not page:
            continue
            
        lines = [line.strip() for line in page.split('\n') if line.strip()]
        if not lines:
            continue
            
        title_text = lines[0].replace("[", "").replace("]", "")
        bullet_points = [l for l in lines[1:] if l.startswith('-') or l.startswith('•') or re.match(r'^\d+\.', l)]
        
        if not bullet_points and len(lines) > 1:
            bullet_points = lines[1:]

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title_text
        
        if slide.placeholders[1].has_text_frame:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for bp in bullet_points:
                p = tf.add_paragraph()
                p.text = bp.lstrip('- •').strip()
                p.level = 0
        
        slide_count += 1

    if slide_count == 0:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "重點摘要"
        slide.placeholders[1].text = content

    # 5. 存檔到相同資料夾
    output_name = file_path.with_name(file_path.stem + "_重點簡報.pptx")
    try:
        prs.save(str(output_name))
        print(f"完成！投影片已存至: {output_name}")
        
        # 寄送通知 (帶上文字摘要)
        send_notification(str(file_path), str(output_name), content)
        
    except Exception as e:
        print(f"存檔失敗: {e}")

if __name__ == "__main__":
    if len(sys.argv) > 1:
        process_media_to_slides(sys.argv[1])
    else:
        print("請提供檔案路徑作為參數。")
